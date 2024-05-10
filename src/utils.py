import psycopg2
from psycopg2.extras import Json
from typing import Optional, Dict, List
import os
import logging
import re
from io import BytesIO
from pdfminer.high_level import extract_text
from docx import Document
from pptx import Presentation
from bs4 import BeautifulSoup
import requests
from cachetools import TTLCache
from openai import AsyncOpenAI
import tempfile

# Load environment variables
from dotenv import load_dotenv
load_dotenv()

# Database connection
database_url = os.getenv("DATABASE_URL")
slack_bot_token = os.getenv("SLACK_BOT_TOKEN")
processed_events = TTLCache(maxsize=1000, ttl=60)

# Database connection function
def get_db_connection(autocommit=True):
    conn = psycopg2.connect(database_url)
    if autocommit:
        conn.autocommit = True
    return conn

def store_transcript(conversation_id, user_id, channel_id, thread_ts, title, transcript):
    with get_db_connection() as conn:
        with conn.cursor() as cur:
            cur.execute(
                """
                INSERT INTO transcripts (conversation_id, user_id, channel_id, thread_ts, title, transcript)
                VALUES (%s, %s, %s, %s, %s, %s)
                ON CONFLICT (conversation_id)
                DO UPDATE SET title = EXCLUDED.title, transcript = EXCLUDED.transcript, created_at = NOW();
                """,
                (conversation_id, user_id, channel_id, thread_ts, title, transcript)
            )

def get_transcript(title):
    with get_db_connection() as conn:
        with conn.cursor() as cur:
            cur.execute(
                "SELECT transcript FROM transcripts WHERE title = %s",
                (title,)
            )
            result = cur.fetchone()
            return result[0] if result else None

def create_message_blocks(text_responses: List[str], button_payloads: Dict) -> (List[Dict], str):
    blocks = []
    summary_text = "Select an option:"
    max_chars = 3000  # Maximum characters for a block of text

    def split_text(text, max_length):
        for start in range(0, len(text), max_length):
            yield text[start:start + max_length]

    for text in text_responses:
        if len(text) <= max_chars:
            blocks.append({"type": "divider"})
            blocks.append({
                "type": "section",
                "text": {
                    "type": "mrkdwn",
                    "text": text
                }
            })
        else:
            for chunk in split_text(text, max_chars):
                blocks.append({"type": "divider"})
                blocks.append({
                    "type": "section",
                    "text": {
                        "type": "mrkdwn",
                        "text": chunk
                    }
                })

    blocks.append({"type": "divider"})
    buttons = []
    for idx, (button_value, button_payload) in enumerate(button_payloads.items()):
        button_text = button_payload['payload']['label']
        buttons.append({
            "type": "button",
            "text": {
                "type": "plain_text",
                "text": button_text,
                "emoji": True
            },
            "value": button_value,
            "action_id": f"voiceflow_button_{idx}"
        })

    if buttons:
        blocks.append({
            "type": "actions",
            "elements": buttons
        })

    return blocks, summary_text

async def download_file(file_url):
    headers = {'Authorization': f'Bearer {slack_bot_token}'}
    response = requests.get(file_url, headers=headers, allow_redirects=True)
    if response.status_code == 200:
        file_suffix = ".mp4" if file_url.endswith(".mp4") else ".m4a" if file_url.endswith(".m4a") else ""
        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=file_suffix)
        temp_file.write(response.content)
        temp_file.close()
        return temp_file.name
    else:
        logging.error(f"Error downloading file: {response.status_code}, {response.text}")
        return None

async def process_file(file_url, file_type):
    file_path = await download_file(file_url)
    if not file_path:
        return None

    try:
        if file_type in ['mp4', 'm4a']:  # Check for both mp4 and m4a file types
            logging.info(f"File path: {file_path}")
            logging.info(f"File size: {os.path.getsize(file_path)}")

            with open(file_path, "rb") as file_stream:
                transcription = await transcribe_audio(file_stream)
                if transcription is None:
                    logging.error("Failed to transcribe or no transcription returned")
                    return None
                else:
                    return create_text_file_in_memory(transcription)
        elif file_type == 'pdf':
            return extract_text_from_pdf(file_content)
        elif file_type in ['doc', 'docx']:
            return extract_text_from_docx(file_content)
        elif file_type in ['ppt', 'pptx']:
            return extract_text_from_pptx(file_content)
    except Exception as e:
        logging.error(f"General error processing file: {e}")
    finally:
        if os.path.exists(file_path):
            os.unlink(file_path)

def extract_text_from_pdf(file_content):
    try:
        file_stream = BytesIO(file_content)
        text = extract_text(file_stream)
        return text
    except Exception as e:
        logging.error(f"Error extracting text from PDF: {e}")
        return None

def extract_text_from_docx(file_content):
    try:
        file_stream = BytesIO(file_content)
        doc = Document(file_stream)
        return "\n".join([paragraph.text for paragraph in doc.paragraphs])
    except Exception as e:
        logging.error(f"Error extracting text from DOCX: {e}")
        return None

def extract_text_from_pptx(file_content):
    try:
        file_stream = BytesIO(file_content)
        ppt = Presentation(file_stream)
        text = []
        for slide in ppt.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)
    except Exception as e:
        logging.error(f"Error extracting text from PPTX: {e}")
        return None

def extract_webpage_content(url):
    headers = {
        'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36',
        'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,image/webp,*/*;q=0.8',
        'Accept-Language': 'en-US,en;q=0.5',
        'Referer': 'https://www.google.com/'
    }
    try:
        response = requests.get(url, headers=headers)
        response.raise_for_status()
        soup = BeautifulSoup(response.content, 'html.parser')
        content_list = []
        for tag in soup.find_all(['p', 'h1', 'h2', 'h3', 'h4', 'h5', 'h6']):
            tag_text = tag.get_text(separator=" ", strip=True)
            clean_text = re.sub(r'\s+', ' ', tag_text).strip()
            if clean_text:
                content_list.append((tag.name, clean_text))
        full_text = ' '.join([text for _, text in content_list])
        return full_text
    except requests.HTTPError as http_err:
        logging.error(f"HTTP error occurred while fetching content from {url}: {http_err}")
        return None
    except Exception as e:
        logging.error(f"An error occurred while fetching content from {url}: {e}")
        return None

async def transcribe_audio(file_stream):
    try:
        openai_api_key = os.getenv("OPENAI_API_KEY")
        openai_client = AsyncOpenAI(api_key=openai_api_key)

        logging.info("Making API call to transcribe audio")
        transcription_response = await openai_client.audio.transcriptions.create(
            model="whisper-1",
            file=file_stream,
            response_format="vtt",
            language="en"
        )
        return transcription_response
    except Exception as e:
        logging.error(f"Error during transcription: {str(e)}")
        return None

def create_text_file_in_memory(content):
    if content is None:
        logging.error("No content to encode into memory")
        return None
    text_stream = BytesIO(content.encode('utf-8'))
    text_stream.seek(0)
    return text_stream